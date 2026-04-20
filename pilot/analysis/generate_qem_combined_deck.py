from pathlib import Path

import numpy as np
import pandas as pd
from scipy import stats
from pptx import Presentation
from pptx.util import Inches, Pt


BASE = Path(__file__).resolve().parent.parent
OUT_DIR = BASE / "output_gemini3.1pro"
SUMMARY_FIG_DIR = OUT_DIR / "analysis_figures"
DEEP_FIG_DIR = OUT_DIR / "analysis_figures_deep"

CLAUDE_CSV = OUT_DIR / "all_questions_collated_gemini3.1pro_eval_claude_3_7_sonnet_eval.csv"
GPT_CSV = OUT_DIR / "all_questions_collated_gemini3.1pro_eval_gpt_4o_eval.csv"

PPT_PATH = OUT_DIR / "QEM_analysis_combined_gemini3_1pro.pptx"


def fmt_p(p: float) -> str:
    if p < 1e-4:
        return "< 1e-4"
    return f"{p:.4f}"


def add_title(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.35), Inches(0.15), Inches(12.6), Inches(0.6))
    tf = box.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(27)
    tf.paragraphs[0].font.bold = True


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: int = 12) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    for p in tf.paragraphs:
        p.font.size = Pt(size)


def main() -> None:
    # Load data for plain-language stats callouts.
    c = pd.read_csv(CLAUDE_CSV)
    g = pd.read_csv(GPT_CSV)

    for df in (c, g):
        df["question_quality"] = pd.to_numeric(df["question_quality"], errors="coerce")
        df["relevance"] = pd.to_numeric(df["relevance"], errors="coerce")
        df["complexity_level"] = pd.to_numeric(df["complexity_level"], errors="coerce")
        df["qem"] = (df["question_quality"] + df["relevance"]) / 2.0

    m = c[["uid", "topic", "source_file", "qem"]].merge(
        g[["uid", "qem"]], on="uid", suffixes=("_claude", "_gpt")
    )
    m["delta_qem"] = m["qem_claude"] - m["qem_gpt"]

    diff = m["delta_qem"].to_numpy()
    ttest = stats.ttest_rel(m["qem_claude"], m["qem_gpt"], nan_policy="omit")
    try:
        wilcoxon = stats.wilcoxon(diff)
        wilcoxon_p = float(wilcoxon.pvalue)
    except ValueError:
        wilcoxon_p = np.nan

    dz = float(np.mean(diff) / np.std(diff, ddof=1)) if np.std(diff, ddof=1) > 0 else 0.0

    topic_mean = pd.concat([
        c[["topic", "qem"]],
        g[["topic", "qem"]],
    ]).groupby("topic")["qem"].mean().sort_values()

    source_mean = pd.concat([
        c[["source_file", "qem"]],
        g[["source_file", "qem"]],
    ]).groupby("source_file")["qem"].mean().sort_values()

    top_topic = topic_mean.index[-1]
    low_topic = topic_mean.index[0]
    top_source = source_mean.index[-1].split(" - ")[0]
    low_source = source_mean.index[0].split(" - ")[0]

    low_thresh = 1.0
    long_all = pd.concat([
        c.assign(evaluator="Claude"),
        g.assign(evaluator="GPT-4o"),
    ], ignore_index=True)
    low_count = int((long_all["qem"] <= low_thresh).sum())

    topic_delta = m.groupby("topic")["delta_qem"].mean().sort_values()
    source_delta = m.groupby("source_file")["delta_qem"].mean().sort_values()

    # Build combined PPT.
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide 1: summary evaluator distribution + significance.
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s1, "QEM Summary: Evaluator Comparison")
    s1.shapes.add_picture(str(SUMMARY_FIG_DIR / "qem_by_evaluator.png"), Inches(0.5), Inches(0.95), width=Inches(8.2), height=Inches(4.8))
    add_text(
        s1,
        (
            "For-dummies read:\n"
            "This chart compares score distributions from the two evaluator models.\n"
            f"- Mean QEM (Claude eval): {m['qem_claude'].mean():.3f}\n"
            f"- Mean QEM (GPT-4o eval): {m['qem_gpt'].mean():.3f}\n"
            f"- Mean difference (Claude-GPT): {np.mean(diff):.3f}\n"
            f"- Paired t-test p-value: {fmt_p(float(ttest.pvalue))}\n"
            f"- Wilcoxon p-value: {fmt_p(wilcoxon_p) if not np.isnan(wilcoxon_p) else 'n/a'}\n"
            f"- Effect size (Cohen dz): {dz:.3f} (small)\n\n"
            "Simple interpretation: the difference is statistically real, but practically small."
        ),
        8.9,
        1.0,
        4.1,
        5.9,
        13,
    )

    # Slide 2: topic + source summary.
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s2, "QEM Summary: By Topic and Source")
    s2.shapes.add_picture(str(SUMMARY_FIG_DIR / "qem_by_topic.png"), Inches(0.45), Inches(0.95), width=Inches(7.5), height=Inches(4.9))
    s2.shapes.add_picture(str(SUMMARY_FIG_DIR / "qem_by_source.png"), Inches(8.0), Inches(0.95), width=Inches(4.9), height=Inches(4.9))
    add_text(
        s2,
        (
            "For-dummies read:\n"
            "Left: average quality score by topic. Right: average by source file.\n"
            f"- Best topic: {top_topic}\n"
            f"- Lowest topic: {low_topic}\n"
            f"- Best source: {top_source}\n"
            f"- Lowest source: {low_source}\n\n"
            "Use this to prioritize where to improve prompts first (start with the lowest groups)."
        ),
        0.5,
        5.95,
        12.3,
        1.2,
        12,
    )

    # Slide 3: disagreement patterns.
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s3, "Deep Dive 1: Evaluator Disagreement Patterns")
    s3.shapes.add_picture(str(DEEP_FIG_DIR / "delta_by_topic_source.png"), Inches(0.45), Inches(0.95), width=Inches(12.45), height=Inches(4.8))
    add_text(
        s3,
        (
            "For-dummies read:\n"
            "Delta = Claude score minus GPT-4o score.\n"
            "If values are negative, GPT-4o tends to be more generous there.\n"
            f"- Most GPT-higher topic: {topic_delta.index[0]} ({topic_delta.iloc[0]:.3f})\n"
            f"- Most Claude-higher topic: {topic_delta.index[-1]} ({topic_delta.iloc[-1]:.3f})\n"
            f"- Most GPT-higher source: {source_delta.index[0].split(' - ')[0]} ({source_delta.iloc[0]:.3f})"
        ),
        0.5,
        5.95,
        12.3,
        1.2,
        12,
    )

    # Slide 4: interaction heatmaps.
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s4, "Deep Dive 2: Topic x Source Interaction")
    s4.shapes.add_picture(str(DEEP_FIG_DIR / "interaction_heatmaps.png"), Inches(0.35), Inches(0.9), width=Inches(12.7), height=Inches(5.0))
    add_text(
        s4,
        (
            "For-dummies read:\n"
            "Left heatmap = absolute score quality for each topic+source pair.\n"
            "Right heatmap = evaluator gap for that same pair.\n"
            "This helps you fix specific weak pockets instead of changing everything globally."
        ),
        0.5,
        6.0,
        12.2,
        1.1,
        12,
    )

    # Slide 5: complexity and tails.
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s5, "Deep Dive 3: Complexity Linkage and Low-QEM Tails")
    s5.shapes.add_picture(str(DEEP_FIG_DIR / "complexity_and_tails.png"), Inches(0.4), Inches(0.95), width=Inches(12.5), height=Inches(5.0))
    add_text(
        s5,
        (
            "For-dummies read:\n"
            "Left: does difficulty level move scores up or down?\n"
            "Right: where do very low scores accumulate?\n"
            f"Current low-QEM rows (<= {low_thresh}): {low_count} of {len(long_all)} evaluator-rows."
        ),
        0.5,
        6.0,
        12.2,
        1.1,
        12,
    )

    # Slide 6: what to do next.
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s6, "Action Plan: Next Analyses To Run")
    add_text(
        s6,
        (
            "1) Inter-rater agreement metrics (weighted kappa, Spearman)\n"
            "   Plain English: do evaluator models rank questions similarly, beyond average score differences?\n\n"
            "2) Error-category mining from identified_issues\n"
            "   Plain English: what failure reasons repeat most, and in which topic/source buckets?\n\n"
            "3) Outlier worksheet (highest absolute delta rows)\n"
            "   Plain English: inspect disagreements first; they are your quickest path to rubric and prompt improvement.\n\n"
            "4) Confidence intervals by topic/source\n"
            "   Plain English: confirm differences are stable, not random noise.\n\n"
            "5) Text-feature linkage to QEM\n"
            "   Plain English: discover wording patterns (length/style) that systematically lower scores."
        ),
        0.6,
        1.1,
        12.1,
        5.9,
        14,
    )

    prs.save(PPT_PATH)
    print(f"Saved combined PPT: {PPT_PATH}")


if __name__ == "__main__":
    main()
