from pathlib import Path

import numpy as np
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import seaborn as sns
from scipy import stats
from pptx import Presentation
from pptx.util import Inches, Pt


BASE = Path(__file__).resolve().parent.parent
OUT_DIR = BASE / "output_gemini3.1pro"
CLAUDE_CSV = OUT_DIR / "all_questions_collated_gemini3.1pro_eval_claude_3_7_sonnet_eval.csv"
GPT_CSV = OUT_DIR / "all_questions_collated_gemini3.1pro_eval_gpt_4o_eval.csv"
FIG_DIR = OUT_DIR / "analysis_figures"
PPT_PATH = OUT_DIR / "QEM_analysis_gemini3_1pro.pptx"
SUMMARY_PATH = OUT_DIR / "QEM_analysis_summary.csv"


def load_eval(path: Path, evaluator_name: str) -> pd.DataFrame:
    df = pd.read_csv(path)
    for col in ["question_quality", "relevance", "complexity_level"]:
        df[col] = pd.to_numeric(df[col], errors="coerce")

    df["qem"] = (df["question_quality"] + df["relevance"]) / 2.0
    df["evaluator"] = evaluator_name
    return df


def cohen_dz(diff: np.ndarray) -> float:
    sd = np.std(diff, ddof=1)
    if sd == 0:
        return 0.0
    return float(np.mean(diff) / sd)


def fmt_p(p: float) -> str:
    if p < 1e-4:
        return "< 1e-4"
    return f"{p:.4f}"


def main() -> None:
    FIG_DIR.mkdir(parents=True, exist_ok=True)

    claude = load_eval(CLAUDE_CSV, "Claude-3.7-Sonnet evaluator")
    gpt = load_eval(GPT_CSV, "GPT-4o evaluator")

    merged = claude[["uid", "qem", "topic", "source_file"]].merge(
        gpt[["uid", "qem"]], on="uid", suffixes=("_claude", "_gpt")
    )

    # Significance tests on paired QEM scores for the same questions.
    diff = (merged["qem_claude"] - merged["qem_gpt"]).to_numpy()
    t_stat, t_p = stats.ttest_rel(merged["qem_claude"], merged["qem_gpt"], nan_policy="omit")
    try:
        w_stat, w_p = stats.wilcoxon(diff)
    except ValueError:
        w_stat, w_p = np.nan, np.nan

    dz = cohen_dz(diff)

    long_df = pd.concat([claude, gpt], ignore_index=True)

    # Aggregates
    model_summary = long_df.groupby("evaluator")["qem"].agg(["count", "mean", "median", "std"]).reset_index()
    topic_summary = long_df.groupby("topic")["qem"].agg(["count", "mean", "median", "std"]).sort_values("mean", ascending=False).reset_index()
    source_summary = long_df.groupby("source_file")["qem"].agg(["count", "mean", "median", "std"]).sort_values("mean", ascending=False).reset_index()

    # Save summary table
    model_summary.to_csv(SUMMARY_PATH, index=False)

    sns.set_theme(style="whitegrid")

    # Figure 1: evaluator distributions
    fig1, ax1 = plt.subplots(figsize=(8.2, 4.6))
    sns.violinplot(data=long_df, x="evaluator", y="qem", inner="box", cut=0, ax=ax1, palette="Set2")
    ax1.set_title("QEM Distribution by Evaluating Model")
    ax1.set_xlabel("")
    ax1.set_ylabel("QEM = (question_quality + relevance) / 2")
    fig1.tight_layout()
    fig1_path = FIG_DIR / "qem_by_evaluator.png"
    fig1.savefig(fig1_path, dpi=220)
    plt.close(fig1)

    # Figure 2: topic distribution (mean + CI bars)
    fig2, ax2 = plt.subplots(figsize=(10.2, 5.6))
    topic_plot = topic_summary.copy()
    sns.barplot(data=topic_plot, x="mean", y="topic", ax=ax2, color="#4c78a8")
    ax2.set_title("Mean QEM by Topic (Both Evaluators Combined)")
    ax2.set_xlabel("Mean QEM")
    ax2.set_ylabel("Topic")
    fig2.tight_layout()
    fig2_path = FIG_DIR / "qem_by_topic.png"
    fig2.savefig(fig2_path, dpi=220)
    plt.close(fig2)

    # Figure 3: source distribution
    fig3, ax3 = plt.subplots(figsize=(9.2, 4.8))
    src = source_summary.copy()
    src["source_short"] = src["source_file"].str.replace("_q_openrouter.csv", "", regex=False)
    sns.barplot(data=src, x="mean", y="source_short", ax=ax3, color="#59a14f")
    ax3.set_title("Mean QEM by Source File (Both Evaluators Combined)")
    ax3.set_xlabel("Mean QEM")
    ax3.set_ylabel("Source file")
    fig3.tight_layout()
    fig3_path = FIG_DIR / "qem_by_source.png"
    fig3.savefig(fig3_path, dpi=220)
    plt.close(fig3)

    # Build PPT
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide 1
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = s1.shapes.add_textbox(Inches(0.35), Inches(0.15), Inches(12.6), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "QEM Overview: Evaluator Agreement and Differences"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True

    s1.shapes.add_picture(str(fig1_path), Inches(0.5), Inches(0.95), width=Inches(8.1))

    stats_box = s1.shapes.add_textbox(Inches(8.8), Inches(1.05), Inches(4.2), Inches(5.8))
    st = stats_box.text_frame
    st.word_wrap = True

    m_cl = model_summary.loc[model_summary["evaluator"].str.contains("Claude"), "mean"].iloc[0]
    m_gp = model_summary.loc[model_summary["evaluator"].str.contains("GPT"), "mean"].iloc[0]

    lines = [
        "Significance (paired, n=570)",
        f"- Mean QEM (Claude eval): {m_cl:.3f}",
        f"- Mean QEM (GPT-4o eval): {m_gp:.3f}",
        f"- Mean diff (Claude-GPT): {np.mean(diff):.3f}",
        f"- Paired t-test p-value: {fmt_p(float(t_p))}",
        f"- Wilcoxon p-value: {fmt_p(float(w_p)) if not np.isnan(w_p) else 'n/a'}",
        f"- Effect size (Cohen dz): {dz:.3f}",
        "",
        "Interpretation",
    ]

    if float(t_p) < 0.05:
        lines.append("- Statistically significant evaluator difference")
    else:
        lines.append("- No statistically significant evaluator difference")

    if abs(dz) < 0.2:
        lines.append("- Practical effect size is tiny")
    elif abs(dz) < 0.5:
        lines.append("- Practical effect size is small")
    else:
        lines.append("- Practical effect size is moderate+")

    st.text = "\n".join(lines)
    for p in st.paragraphs:
        p.font.size = Pt(16)

    # Slide 2
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    t2 = s2.shapes.add_textbox(Inches(0.35), Inches(0.15), Inches(12.6), Inches(0.6))
    t2f = t2.text_frame
    t2f.text = "QEM by Topic and Source File"
    t2f.paragraphs[0].font.size = Pt(28)
    t2f.paragraphs[0].font.bold = True

    s2.shapes.add_picture(str(fig2_path), Inches(0.45), Inches(0.95), width=Inches(7.4))
    s2.shapes.add_picture(str(fig3_path), Inches(7.95), Inches(0.95), width=Inches(5.0))

    # Extra insights text
    top_topic = topic_summary.iloc[0]
    bottom_topic = topic_summary.iloc[-1]
    top_source = source_summary.iloc[0]
    bottom_source = source_summary.iloc[-1]

    notes = s2.shapes.add_textbox(Inches(7.95), Inches(5.35), Inches(5.0), Inches(1.9))
    nt = notes.text_frame
    nt.word_wrap = True
    nt.text = (
        "What else to look at next:\n"
        f"- Top topic by mean QEM: {top_topic['topic']} ({top_topic['mean']:.2f})\n"
        f"- Lowest topic by mean QEM: {bottom_topic['topic']} ({bottom_topic['mean']:.2f})\n"
        f"- Top source: {str(top_source['source_file']).split(' - ')[0]} ({top_source['mean']:.2f})\n"
        f"- Lowest source: {str(bottom_source['source_file']).split(' - ')[0]} ({bottom_source['mean']:.2f})\n"
        "- Also examine evaluator disagreement by topic/source"
    )
    for p in nt.paragraphs:
        p.font.size = Pt(13)

    prs.save(PPT_PATH)

    print(f"Saved PPT: {PPT_PATH}")
    print(f"Saved summary: {SUMMARY_PATH}")
    print("Model means:")
    print(model_summary.to_string(index=False))


if __name__ == "__main__":
    main()
