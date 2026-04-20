from pathlib import Path

import numpy as np
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt


BASE = Path(__file__).resolve().parent.parent
OUT_DIR = BASE / "output_gemini3.1pro"
CLAUDE_CSV = OUT_DIR / "all_questions_collated_gemini3.1pro_eval_claude_3_7_sonnet_eval.csv"
GPT_CSV = OUT_DIR / "all_questions_collated_gemini3.1pro_eval_gpt_4o_eval.csv"
FIG_DIR = OUT_DIR / "analysis_figures_deep"
PPT_PATH = OUT_DIR / "QEM_analysis_deep_dive_gemini3_1pro.pptx"


def add_title(slide, text: str) -> None:
    tb = slide.shapes.add_textbox(Inches(0.35), Inches(0.15), Inches(12.6), Inches(0.6))
    tf = tb.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(27)
    tf.paragraphs[0].font.bold = True


def add_explainer(slide, x: float, y: float, w: float, h: float, text: str, size: int = 12) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    for p in tf.paragraphs:
        p.font.size = Pt(size)


def main() -> None:
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    sns.set_theme(style="whitegrid")

    c = pd.read_csv(CLAUDE_CSV)
    g = pd.read_csv(GPT_CSV)

    for df in (c, g):
        df["question_quality"] = pd.to_numeric(df["question_quality"], errors="coerce")
        df["relevance"] = pd.to_numeric(df["relevance"], errors="coerce")
        df["complexity_level"] = pd.to_numeric(df["complexity_level"], errors="coerce")
        df["qem"] = (df["question_quality"] + df["relevance"]) / 2.0

    # Merge by question uid for paired analyses.
    m = c[["uid", "topic", "source_file", "complexity_level", "qem"]].merge(
        g[["uid", "qem", "complexity_level"]],
        on="uid",
        suffixes=("_claude", "_gpt"),
    )
    m["delta_qem"] = m["qem_claude"] - m["qem_gpt"]

    top_topic_neg = m.groupby("topic")["delta_qem"].mean().sort_values().head(1)
    top_topic_pos = m.groupby("topic")["delta_qem"].mean().sort_values(ascending=False).head(1)
    top_source_neg = m.groupby("source_file")["delta_qem"].mean().sort_values().head(1)

    # 1) Evaluator disagreement by topic and source.
    topic_delta = m.groupby("topic")["delta_qem"].mean().sort_values()
    source_delta = m.groupby("source_file")["delta_qem"].mean().sort_values()

    fig1, (ax1, ax2) = plt.subplots(1, 2, figsize=(12.2, 5.2), gridspec_kw={"width_ratios": [1.15, 1]})
    topic_delta.plot(kind="barh", ax=ax1, color="#e15759")
    ax1.axvline(0, color="black", linewidth=1)
    ax1.set_title("Mean QEM Delta by Topic (Claude - GPT)")
    ax1.set_xlabel("Delta QEM")
    ax1.set_ylabel("Topic")

    source_short_idx = [s.split(" - ")[0] for s in source_delta.index]
    ax2.barh(source_short_idx, source_delta.values, color="#f28e2b")
    ax2.axvline(0, color="black", linewidth=1)
    ax2.set_title("Mean QEM Delta by Source")
    ax2.set_xlabel("Delta QEM")
    ax2.set_ylabel("Source")

    fig1.tight_layout()
    fig1_path = FIG_DIR / "delta_by_topic_source.png"
    fig1.savefig(fig1_path, dpi=220)
    plt.close(fig1)

    # 2) Topic x Source interaction heatmaps.
    long_df = pd.concat([
        c.assign(evaluator="Claude")[ ["topic", "source_file", "qem", "evaluator"] ],
        g.assign(evaluator="GPT-4o")[ ["topic", "source_file", "qem", "evaluator"] ],
    ], ignore_index=True)

    pivot_mean = long_df.pivot_table(index="topic", columns="source_file", values="qem", aggfunc="mean")
    pivot_delta = m.pivot_table(index="topic", columns="source_file", values="delta_qem", aggfunc="mean")

    fig2, (ax3, ax4) = plt.subplots(1, 2, figsize=(12.8, 5.8))
    sns.heatmap(pivot_mean, cmap="YlGnBu", ax=ax3)
    ax3.set_title("Mean QEM: Topic x Source")
    ax3.set_xlabel("Source file")
    ax3.set_ylabel("Topic")

    sns.heatmap(pivot_delta, cmap="RdBu_r", center=0, ax=ax4)
    ax4.set_title("Evaluator Delta (Claude-GPT): Topic x Source")
    ax4.set_xlabel("Source file")
    ax4.set_ylabel("Topic")

    fig2.tight_layout()
    fig2_path = FIG_DIR / "interaction_heatmaps.png"
    fig2.savefig(fig2_path, dpi=220)
    plt.close(fig2)

    # 3) Complexity linkage + tails.
    long_for_complex = pd.concat([
        c.assign(evaluator="Claude")[ ["uid", "topic", "source_file", "complexity_level", "qem", "evaluator"] ],
        g.assign(evaluator="GPT-4o")[ ["uid", "topic", "source_file", "complexity_level", "qem", "evaluator"] ],
    ], ignore_index=True)

    comp_summary = long_for_complex.groupby(["evaluator", "complexity_level"])["qem"].mean().reset_index()

    # Tail definition: low-qem rows.
    low_thresh = 1.0
    low_rows = long_for_complex[long_for_complex["qem"] <= low_thresh]
    low_topic = low_rows.groupby("topic").size().sort_values(ascending=False).head(10)
    low_source = low_rows.groupby("source_file").size().sort_values(ascending=False)

    fig3, (ax5, ax6) = plt.subplots(1, 2, figsize=(12.6, 5.2))
    sns.lineplot(data=comp_summary, x="complexity_level", y="qem", hue="evaluator", marker="o", ax=ax5)
    ax5.set_title("QEM vs Complexity Level")
    ax5.set_xlabel("Complexity level")
    ax5.set_ylabel("Mean QEM")
    ax5.set_xticks(sorted(comp_summary["complexity_level"].dropna().unique()))

    src_names = [s.split(" - ")[0] for s in low_source.index]
    ax6.barh(src_names, low_source.values, color="#76b7b2")
    ax6.set_title("Low-QEM Count (<= 1.0) by Source")
    ax6.set_xlabel("Count")
    ax6.set_ylabel("Source")

    fig3.tight_layout()
    fig3_path = FIG_DIR / "complexity_and_tails.png"
    fig3.savefig(fig3_path, dpi=220)
    plt.close(fig3)

    # Build slides.
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s1, "Deep Dive 1: Evaluator Disagreement Patterns")
    s1.shapes.add_picture(str(fig1_path), Inches(0.45), Inches(0.95), width=Inches(12.45), height=Inches(4.75))

    add_explainer(
        s1,
        0.5,
        5.85,
        12.2,
        1.35,
        (
            "For-dummies read: each bar shows who is stricter. "
            "Delta = Claude minus GPT-4o. If a bar is below 0, GPT-4o usually gives higher scores there. "
            f"Biggest GPT-higher topic: {top_topic_neg.index[0]} ({top_topic_neg.iloc[0]:.3f}). "
            f"Closest to Claude-higher topic: {top_topic_pos.index[0]} ({top_topic_pos.iloc[0]:.3f}). "
            f"Strongest source gap: {top_source_neg.index[0].split(' - ')[0]} ({top_source_neg.iloc[0]:.3f})."
        ),
        size=12,
    )

    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s2, "Deep Dive 2: Topic x Source Interaction")
    s2.shapes.add_picture(str(fig2_path), Inches(0.35), Inches(0.9), width=Inches(12.7), height=Inches(4.9))
    add_explainer(
        s2,
        0.45,
        5.95,
        12.45,
        1.2,
        (
            "For-dummies read: left heatmap = absolute quality by topic+source (darker means better). "
            "Right heatmap = evaluator disagreement for that same cell. "
            "Use this to find specific pockets that need prompt tuning, instead of changing everything globally."
        ),
        size=12,
    )

    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s3, "Deep Dive 3: Complexity Linkage and Low-QEM Tails")
    s3.shapes.add_picture(str(fig3_path), Inches(0.4), Inches(0.95), width=Inches(12.5), height=Inches(4.9))
    add_explainer(
        s3,
        0.5,
        5.95,
        12.3,
        1.2,
        (
            "For-dummies read: left chart asks 'do harder questions get better/worse scores?'. "
            "Right chart shows where low scores pile up. "
            f"Current low-QEM count (<= {low_thresh}): {len(low_rows)} / {len(long_for_complex)} evaluator-rows."
        ),
        size=12,
    )

    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s4, "What Else To Check Next (Simple Guide)")
    add_explainer(
        s4,
        0.55,
        1.05,
        12.1,
        5.9,
        (
            "1) Inter-rater agreement (weighted kappa, Spearman):\n"
            "   Plain English: do the two evaluator models actually agree on ranking quality, or just have similar averages?\n\n"
            "2) Issue-type mining from identified_issues:\n"
            "   Plain English: what are the most common failure reasons (trivial question, ambiguity, grammar), and where do they happen most?\n\n"
            "3) Outlier review sheet (top disagreement rows):\n"
            "   Plain English: inspect the exact questions where evaluators disagree most; those rows are your fastest path to rubric/prompt fixes.\n\n"
            "4) Topic confidence intervals, not just means:\n"
            "   Plain English: check if differences are stable or just noise before making product decisions.\n\n"
            "5) Text-feature linkage (question length, style, punctuation) vs QEM:\n"
            "   Plain English: detect easy-to-fix wording patterns that systematically lower scores."
        ),
        size=14,
    )

    prs.save(PPT_PATH)
    print(f"Saved deep-dive PPT: {PPT_PATH}")
    print(f"Low-QEM rows (<= {low_thresh}): {len(low_rows)} out of {len(long_for_complex)} evaluator-rows")


if __name__ == "__main__":
    main()
